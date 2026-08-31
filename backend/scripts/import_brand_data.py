"""一次性导入美丽田园品牌资料到知识库。"""

from __future__ import annotations

import hashlib
import html
import os
import re
import subprocess
import sys
from pathlib import Path

from sqlalchemy.orm import Session

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from app.config import get_settings  # noqa: E402
from app.db import SessionLocal  # noqa: E402
from app.knowledge.service import KnowledgeService  # noqa: E402
from app.models import KnowledgeDoc  # noqa: E402

ROOT = Path("/Users/zhaoxinyuan/Desktop/消费者运营中台/美丽田园")
MAX_FILE_BYTES = 8 * 1024 * 1024
SKIP_DIR_PARTS = {
    ".build",
    ".git",
    ".swiftpm",
    "SecondBrain",
    "SQLBrain",
    "SQLBrainApp",
    "SB_HF",
    "Sources",
    "ZOERA",
    "ZOÉRA.app",
    "node_modules",
    "templates",
}
SKIP_EXTS = {".ds_store", ".png", ".xmind", ".app", ".pyc"}


def _iter_files() -> list[tuple[str, Path]]:
    rows: list[tuple[str, Path]] = []
    for path in sorted(ROOT.rglob("*")):
        if not path.is_file():
            continue
        if any(part in SKIP_DIR_PARTS for part in path.parts):
            continue
        if path.suffix.lower() in SKIP_EXTS:
            continue
        rel = path.relative_to(ROOT).as_posix()
        rows.append((rel, path))
    return rows


def _redact(text: str) -> str:
    text = re.sub(r"1[3-9]\d{9}", "[手机号已脱敏]", text)
    text = re.sub(r"wxid_[A-Za-z0-9_-]+", "[微信号已脱敏]", text)
    text = re.sub(r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}", "[邮箱已脱敏]", text)
    return text


def _sheet_text(sheet) -> str:
    lines: list[str] = []
    for row_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
        if row_idx > 300:
            break
        values = ["" if v is None else str(v).strip() for v in row[:40]]
        line = " | ".join(v for v in values if v)
        if line:
            lines.append(line)
    return "\n".join(lines)


def extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets[:8]:
        text = _sheet_text(sheet)
        if text.strip():
            parts.append(f"## 工作表：{sheet.title}\n{text}")
    wb.close()
    return "\n\n".join(parts)


def extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables[:20]:
        for row in table.rows[:100]:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            line = " | ".join(c for c in cells if c)
            if line:
                parts.append(line)
    return "\n".join(parts)


def extract_doc(path: Path) -> str:
    result = subprocess.run(
        ["textutil", "-convert", "txt", "-stdout", str(path)],
        capture_output=True,
        text=True,
        timeout=60,
        check=False,
    )
    return result.stdout or ""


def extract_pdf(path: Path) -> str:
    import pdfplumber

    parts: list[str] = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages[:80]:
            text = page.extract_text() or ""
            if text.strip():
                parts.append(text)
    return "\n\n".join(parts)


def extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    parts: list[str] = []
    for slide in prs.slides[:60]:
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text.strip())
        if texts:
            parts.append(" / ".join(t for t in texts if t))
    return "\n\n".join(parts)


def extract_text_file(path: Path) -> str:
    content = path.read_text(encoding="utf-8", errors="replace")
    if path.suffix.lower() in {".html", ".htm"}:
        content = re.sub(r"<script.*?</script>", " ", content, flags=re.S | re.I)
        content = re.sub(r"<style.*?</style>", " ", content, flags=re.S | re.I)
        content = re.sub(r"<[^>]+>", " ", content)
        content = html.unescape(content)
    return content


def extract(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".xlsx":
        return extract_xlsx(path)
    if ext == ".docx":
        return extract_docx(path)
    if ext == ".doc":
        return extract_doc(path)
    if ext == ".pdf":
        return extract_pdf(path)
    if ext == ".pptx":
        return extract_pptx(path)
    if ext in {".txt", ".md", ".csv", ".json", ".html", ".htm"}:
        return extract_text_file(path)
    return ""


def main() -> None:
    settings = get_settings()
    files = _iter_files()
    print(f"共发现 {len(files)} 个待处理文件")

    with SessionLocal() as db:
        service = KnowledgeService(db)
        existing = {
            row.name
            for row in db.query(KnowledgeDoc).filter(KnowledgeDoc.tenant_id == settings.default_tenant_id).all()
        }
        seen_hashes: set[str] = set()
        imported = 0
        skipped_size = 0
        skipped_empty = 0
        skipped_dup = 0
        failed: list[str] = []

        for rel, path in files:
            size = path.stat().st_size
            if size > MAX_FILE_BYTES:
                print(f"[跳过-超大] {rel} ({size / 1024 / 1024:.1f} MB)")
                skipped_size += 1
                continue
            try:
                text = extract(path)
            except Exception as exc:  # noqa: BLE001
                failed.append(f"{rel}: {exc}")
                print(f"[失败] {rel}: {exc}")
                continue

            text = _redact(text).strip()
            if len(text) < 30:
                print(f"[跳过-低内容] {rel} ({len(text)} 字)")
                skipped_empty += 1
                continue

            digest = hashlib.md5(text.encode("utf-8")).hexdigest()
            if digest in seen_hashes:
                print(f"[跳过-重复] {rel}")
                skipped_dup += 1
                continue
            seen_hashes.add(digest)

            clean_rel = rel.replace(os.sep, "-").replace(" ", "")
            name = f"美丽田园-{clean_rel}"[:200]
            if name in existing:
                print(f"[跳过-已存在] {rel}")
                skipped_dup += 1
                continue

            doc = service.ingest(
                settings.default_tenant_id,
                name,
                text,
                path.suffix.lower().lstrip(".") or "txt",
                size,
            )
            existing.add(name)
            imported += 1
            print(f"[导入] {rel} -> {doc.chunk_count} 切片")

        print("\n===== 导入汇总 =====")
        print(f"导入成功: {imported}")
        print(f"跳过超大文件: {skipped_size}")
        print(f"跳过低内容/图片型PDF: {skipped_empty}")
        print(f"跳过重复/已存在: {skipped_dup}")
        print(f"失败: {len(failed)}")
        for item in failed:
            print(f"  - {item}")


if __name__ == "__main__":
    main()

