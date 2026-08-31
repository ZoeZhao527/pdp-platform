"""一次性导入华润怡宝品牌资料到知识库（支持超大文件，不跳过）。"""

from __future__ import annotations

import hashlib
import html
import re
import subprocess
import sys
import zipfile
from pathlib import Path
from typing import Iterator
from xml.etree import ElementTree as ET

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from app.db import SessionLocal  # noqa: E402
from app.knowledge.service import KnowledgeService  # noqa: E402
from app.models import KnowledgeDoc  # noqa: E402

ROOT = Path("/Users/zhaoxinyuan/Desktop/消费者运营中台/怡宝/怡宝")
TENANT_ID = "4510b2c8-9761-4b94-ae3e-77b2838906c9"
INDUSTRY_ID = "2a6d6e26-4c2b-43f3-8337-52e05142a726"  # 零售
MAX_ROWS_PER_PART = 2000
MAX_COLS = 80
MAX_FILE_BYTES = 60 * 1024 * 1024  # 60MB — still try, just warn

SKIP_DIR_PARTS = {".DS_Store"}
SKIP_EXTS = {".ds_store", ".png", ".jpg", ".jpeg", ".psd", ".mp4", ".xmind", ".app", ".pyc"}


def _iter_files() -> list[tuple[str, Path]]:
    rows: list[tuple[str, Path]] = []
    for path in sorted(ROOT.rglob("*")):
        if not path.is_file():
            continue
        if path.name == ".DS_Store":
            continue
        if path.suffix.lower() in SKIP_EXTS:
            continue
        rel = path.relative_to(ROOT).as_posix()
        rows.append((rel, path))
    return rows


def _redact(text: str) -> str:
    text = re.sub(r"1[3-9]\d{9}", "[手机号已脱敏]", text)
    text = re.sub(r"wxid_[A-Za-z0-9_-]+", "[微信号已脱敏]", text)
    text = re.sub(r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}", "[邮箱已脱敏]", text)
    return text


def _line(values: list) -> str:
    cleaned = [str(v).strip() for v in values[:MAX_COLS] if v is not None and str(v).strip()]
    return " | ".join(cleaned)


def _parts_from_lines(sheet: str, lines: list[str]) -> Iterator[tuple[str, int, str]]:
    part = 1
    for start in range(0, len(lines), MAX_ROWS_PER_PART):
        chunk = lines[start : start + MAX_ROWS_PER_PART]
        if chunk:
            yield sheet, part, "\n".join(chunk)
        part += 1


def iter_xlsx_parts(path: Path) -> Iterator[tuple[str, int, str]]:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        for ws in wb.worksheets:
            lines = [line for line in (_line(row) for row in ws.iter_rows(values_only=True)) if line]
            yield from _parts_from_lines(ws.title, lines)
        wb.close()
        return
    except Exception:
        pass

    NS = {"m": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
    with zipfile.ZipFile(path) as z:
        shared: list[str] = []
        if "xl/sharedStrings.xml" in z.namelist():
            root = ET.fromstring(z.read("xl/sharedStrings.xml"))
            for si in root.findall("m:si", NS):
                shared.append("".join(t.text or "" for t in si.iter("{http://schemas.openxmlformats.org/spreadsheetml/2006/main}t")))
        wb_xml = ET.fromstring(z.read("xl/workbook.xml"))
        sheets = [s.get("name") for s in wb_xml.find("m:sheets", NS)]
        for idx, sheet in enumerate(sheets, start=1):
            xml_path = f"xl/worksheets/sheet{idx}.xml"
            if xml_path not in z.namelist():
                continue
            root = ET.fromstring(z.read(xml_path))
            lines: list[str] = []
            for row in root.findall(".//m:sheetData/m:row", NS):
                values: list[str] = []
                for cell in row.findall("m:c", NS):
                    cell_type = cell.get("t")
                    value_el = cell.find("m:v", NS)
                    if cell_type == "s" and value_el is not None:
                        values.append(shared[int(value_el.text)])
                    elif cell_type == "inlineStr":
                        inline = cell.find("m:is", NS)
                        values.append(
                            "".join(t.text or "" for t in inline.iter("{http://schemas.openxmlformats.org/spreadsheetml/2006/main}t"))
                            if inline is not None else ""
                        )
                    elif value_el is not None:
                        values.append(value_el.text or "")
                line = _line(values)
                if line:
                    lines.append(line)
            yield from _parts_from_lines(sheet, lines)


def extract_xlsx(path: Path) -> list[tuple[str, str]]:
    """Returns list of (sheet_name, text) for each sheet."""
    parts: list[tuple[str, str]] = []
    for sheet, _part, text in iter_xlsx_parts(path):
        if text.strip():
            parts.append((f"{sheet}", text))
    return parts


def extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables[:20]:
        for row in table.rows[:200]:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            line = " | ".join(c for c in cells if c)
            if line:
                parts.append(line)
    return "\n".join(parts)


def extract_doc(path: Path) -> str:
    result = subprocess.run(
        ["textutil", "-convert", "txt", "-stdout", str(path)],
        capture_output=True, text=True, timeout=60, check=False,
    )
    return result.stdout or ""


def extract_pdf(path: Path) -> str:
    import pdfplumber
    parts: list[str] = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            if text.strip():
                parts.append(text)
    return "\n\n".join(parts)


def extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts: list[str] = []
    for slide in prs.slides[:80]:
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text.strip())
        if texts:
            parts.append(" / ".join(t for t in texts if t))
    return "\n\n".join(parts)


def extract_text_file(path: Path) -> str:
    content = path.read_text(encoding="utf-8", errors="replace")
    if path.suffix.lower() in {".html", ".htm"}:
        content = re.sub(r"<script.*?</script>", " ", content, flags=re.S | re.I)
        content = re.sub(r"<style.*?</style>", " ", content, flags=re.S | re.I)
        content = re.sub(r"<[^>]+>", " ", content)
        content = html.unescape(content)
    return content


def main() -> None:
    files = _iter_files()
    print(f"共发现 {len(files)} 个待处理文件")

    with SessionLocal() as db:
        service = KnowledgeService(db)
        existing = {
            row.name
            for row in db.query(KnowledgeDoc).filter(KnowledgeDoc.tenant_id == TENANT_ID).all()
        }
        seen_hashes: set[str] = set()
        imported = 0
        skipped_empty = 0
        skipped_dup = 0
        failed: list[str] = []

        for rel, path in files:
            size = path.stat().st_size
            ext = path.suffix.lower()
            stem = re.sub(r"[\s]+", "", path.stem)

            if size > MAX_FILE_BYTES:
                print(f"[警告-大文件] {rel} ({size / 1024 / 1024:.1f} MB) 仍尝试导入")

            try:
                if ext == ".xlsx":
                    sheet_texts = extract_xlsx(path)
                    if not sheet_texts:
                        print(f"[跳过-空表] {rel}")
                        skipped_empty += 1
                        continue
                    for sheet_name, text in sheet_texts:
                        text = _redact(text).strip()
                        if len(text) < 10:
                            continue
                        digest = hashlib.md5(text.encode("utf-8")).hexdigest()
                        if digest in seen_hashes:
                            continue
                        seen_hashes.add(digest)
                        name = f"怡宝-{stem}-{sheet_name}.xlsx"[:200]
                        if name in existing:
                            continue
                        doc = service.ingest(
                            TENANT_ID, name, text, "xlsx", size,
                            industry_id=INDUSTRY_ID,
                        )
                        existing.add(name)
                        imported += 1
                        print(f"[导入] {name} -> {doc.chunk_count} 切片")
                    continue

                if ext == ".xls":
                    # Try openpyxl (limited xls support) or textutil
                    try:
                        from openpyxl import load_workbook
                        wb = load_workbook(path, read_only=True, data_only=True)
                        for ws in wb.worksheets:
                            lines = [line for line in (_line(row) for row in ws.iter_rows(values_only=True)) if line]
                            text = _redact("\n".join(lines)).strip()
                            if len(text) < 10:
                                continue
                            digest = hashlib.md5(text.encode("utf-8")).hexdigest()
                            if digest in seen_hashes:
                                continue
                            seen_hashes.add(digest)
                            name = f"怡宝-{stem}-{ws.title}.xls"[:200]
                            if name in existing:
                                continue
                            doc = service.ingest(TENANT_ID, name, text, "xls", size, industry_id=INDUSTRY_ID)
                            existing.add(name)
                            imported += 1
                            print(f"[导入] {name} -> {doc.chunk_count} 切片")
                        wb.close()
                    except Exception:
                        print(f"[跳过-xls解析失败] {rel}")
                        skipped_empty += 1
                    continue

                if ext == ".pdf":
                    text = extract_pdf(path)
                elif ext == ".docx":
                    text = extract_docx(path)
                elif ext == ".doc":
                    text = extract_doc(path)
                elif ext == ".pptx":
                    text = extract_pptx(path)
                elif ext in {".txt", ".md", ".csv", ".json", ".html", ".htm"}:
                    text = extract_text_file(path)
                else:
                    continue

                text = _redact(text).strip()
                if len(text) < 30:
                    print(f"[跳过-低内容] {rel} ({len(text)} 字)")
                    skipped_empty += 1
                    continue

                digest = hashlib.md5(text.encode("utf-8")).hexdigest()
                if digest in seen_hashes:
                    print(f"[跳过-重复] {rel}")
                    skipped_dup += 1
                    continue
                seen_hashes.add(digest)

                name = f"怡宝-{stem}{ext}"[:200]
                if name in existing:
                    print(f"[跳过-已存在] {rel}")
                    skipped_dup += 1
                    continue

                doc = service.ingest(
                    TENANT_ID, name, text, ext.lstrip(".") or "txt", size,
                    industry_id=INDUSTRY_ID,
                )
                existing.add(name)
                imported += 1
                print(f"[导入] {name} -> {doc.chunk_count} 切片（{len(text)} 字）")

            except Exception as exc:
                db.rollback()
                failed.append(f"{rel}: {exc}")
                print(f"[失败] {rel}: {exc}")

        print("\n===== 怡宝导入汇总 =====")
        print(f"导入成功: {imported}")
        print(f"跳过低内容: {skipped_empty}")
        print(f"跳过重复/已存在: {skipped_dup}")
        print(f"失败: {len(failed)}")
        for item in failed:
            print(f"  - {item}")


if __name__ == "__main__":
    main()
